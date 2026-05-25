#!/usr/bin/env python3
"""Generate Cisco Packet Tracer Webinar Script (Word) and Presentation (PowerPoint).
60 Slides - Detailed IP Placement + Add Simple PDU (Mailbox Icon) Testing."""
import os
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

# Ensure output directory exists
os.makedirs("webinar", exist_ok=True)

# Define slide data: each entry is a dict with all info needed
# 60 SLIDES TOTAL
SLIDES = [
    {
        "num": 1,
        "title": "Welcome to Our Cisco Packet Tracer Webinar!",
        "script": (
            "Hello everyone! Welcome to our webinar today. We are so happy you are here. "
            "Today, we will learn about a very cool tool called Cisco Packet Tracer. "
            "It helps us build pretend computer networks on our computer. Think of it like "
            "a video game, but for networking! By the end of this webinar, you will know "
            "exactly how to place IP addresses on every device and test connectivity using "
            "the Add Simple PDU mailbox icon. Let us get started!"
        ),
        "music": "Upbeat instrumental - 'Happy Day' by Bensound",
        "animation": "Fade in with zoom effect on title",
        "emotion": "Excited and energetic - like greeting old friends",
        "bullets": [
            "Welcome!",
            "Topic: Cisco Packet Tracer",
            "Duration: 60 slides, 1+ hour",
            "Focus: IP Addressing + PDU Testing",
            "No terminal/CLI needed!",
            "Interactive session"
        ],
        "bg": "D6EAF8"
    },

    {
        "num": 2,
        "title": "What Will We Learn Today?",
        "script": (
            "Here is what we will talk about today. First, we will learn what Cisco Packet "
            "Tracer is. Then, we will learn how to place IP addresses on PCs using the Desktop "
            "tab and on routers using the Config tab. We will NOT use any terminal or CLI. "
            "Instead, we will use the GUI only. And to test if devices can talk, we will use "
            "the Add Simple PDU icon, which looks like a closed envelope or mailbox on the "
            "right-side toolbar. Sounds fun, right?"
        ),
        "music": "Gentle acoustic guitar - 'New Beginning' by Bensound",
        "animation": "Slide in from left, bullet points appear one by one",
        "emotion": "Warm and inviting - like a teacher starting class",
        "bullets": [
            "What is Cisco Packet Tracer?",
            "How to place IP addresses (GUI only)",
            "PC: Desktop tab > IP Configuration",
            "Router: Config tab > Interface",
            "Test with Add Simple PDU (mailbox icon)",
            "6 Philippine scenarios with full details"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 3,
        "title": "What is Cisco Packet Tracer?",
        "script": (
            "So, what is Cisco Packet Tracer? It is a free program made by Cisco. Cisco is "
            "a very big company that makes network equipment, like routers and switches. "
            "Packet Tracer lets you build a pretend network on your computer. You do not need "
            "real equipment. You can drag and drop devices like computers, routers, and switches. "
            "Then you connect them with cables. It is like building with digital Lego blocks!"
        ),
        "music": "Curious melody - 'Little Idea' by Bensound",
        "animation": "Zoom in on center, text fades in",
        "emotion": "Curious and explanatory - like showing a new toy",
        "bullets": [
            "Free program by Cisco",
            "Build pretend networks",
            "No real equipment needed",
            "Drag and drop devices",
            "Connect with virtual cables",
            "Like digital Lego blocks!"
        ],
        "bg": "FDEBD0"
    },

    {
        "num": 4,
        "title": "What is Cisco Packet Tracer? (continued)",
        "script": (
            "Let me explain more. Packet Tracer is a simulator. That means it pretends to be "
            "a real network. You can see how data moves from one computer to another. You can "
            "watch packets travel through cables and routers. It is used by students all over "
            "the world to learn networking. And the best part? It is completely free if you "
            "sign up for Cisco NetAcad! Many schools in the Philippines use it."
        ),
        "music": "Curious melody continues - 'Little Idea' by Bensound",
        "animation": "Crossfade from previous slide",
        "emotion": "Informative and encouraging - building excitement",
        "bullets": [
            "It is a simulator (pretend network)",
            "Watch data move between computers",
            "See packets travel through cables",
            "Used by students worldwide",
            "Completely FREE with NetAcad account",
            "Many Philippine schools use it"
        ],
        "bg": "FFFFFF"
    },
    {
        "num": 5,
        "title": "Who Made It and Why?",
        "script": (
            "Cisco made Packet Tracer for students. They wanted students to practice networking "
            "without buying expensive equipment. A real Cisco router can cost thousands of pesos! "
            "But with Packet Tracer, you can practice for free. It is part of the Cisco "
            "Networking Academy program. Many schools in the Philippines use it for their IT classes."
        ),
        "music": "Inspiring background - 'Evolution' by Bensound",
        "animation": "Slide in from right with company logo placeholder",
        "emotion": "Appreciative and informative - showing value",
        "bullets": [
            "Made by Cisco for students",
            "Practice without expensive equipment",
            "Real routers cost thousands of pesos",
            "Part of Cisco Networking Academy",
            "Many Philippine schools use it",
            "Perfect for CCNA preparation"
        ],
        "bg": "D6EAF8"
    },

    {
        "num": 6,
        "title": "Key Features - Build Networks Easily",
        "script": (
            "Let us talk about the cool things Packet Tracer can do. First, you can build "
            "networks very easily. Just drag a router from the menu and drop it on the screen. "
            "Then drag a computer. Connect them with a cable. Done! You just built a tiny network. "
            "You can add as many devices as you want. Routers, switches, computers, servers, "
            "phones, even wireless access points!"
        ),
        "music": "Upbeat electronic - 'Jazzy Frenchy' by Bensound",
        "animation": "Pop-in effect for each bullet point",
        "emotion": "Enthusiastic - showing off cool features",
        "bullets": [
            "Drag and drop to build networks",
            "Routers, switches, computers",
            "Servers and phones too",
            "Wireless access points",
            "Connect with virtual cables",
            "Add as many devices as you want"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 7,
        "title": "Key Features - Simulation Mode & Add Simple PDU",
        "script": (
            "The second amazing feature is testing connectivity without typing any commands! "
            "Look at the right side of the Packet Tracer screen. You will see a toolbar with "
            "several icons. One of them looks like a closed envelope or mailbox. This is called "
            "Add Simple PDU. PDU means Protocol Data Unit, which is just a fancy word for a "
            "message. When you click this icon, then click on one device, then click on another "
            "device, Packet Tracer will try to send a message between them. If the message "
            "succeeds, you see Successful in the PDU list at the bottom. If it fails, you see "
            "Failed. No terminal needed!"
        ),
        "music": "Wonder and discovery - 'Creative Minds' by Bensound",
        "animation": "Slide transition with envelope icon highlight",
        "emotion": "Amazed and wonder-filled - like showing magic",
        "bullets": [
            "Right-side toolbar has testing icons",
            "Add Simple PDU = closed envelope/mailbox icon",
            "Click icon, then click source device",
            "Then click destination device",
            "Check bottom panel: Successful or Failed",
            "NO terminal or Command Prompt needed!"
        ],
        "bg": "FDEBD0"
    },

    {
        "num": 8,
        "title": "How to Place IP on a PC (Desktop Tab > IP Configuration)",
        "script": (
            "Now let me teach you the MOST important skill: placing an IP address on a PC. "
            "This is how you do it step by step. Step 1: Click on the PC in your workspace. "
            "A window will pop up showing tabs at the top: Physical, Config, Desktop, etc. "
            "Step 2: Click on the Desktop tab. Step 3: You will see several icons like IP "
            "Configuration, Command Prompt, Web Browser, etc. Click on IP Configuration. "
            "Step 4: You will see fields for IPv4 Address, Subnet Mask, and Default Gateway. "
            "Step 5: Make sure the Static radio button is selected (not DHCP). "
            "Step 6: Type the IP address in the IPv4 Address field. "
            "Step 7: Type the subnet mask (usually 255.255.255.0). "
            "Step 8: Type the default gateway (this is the router IP address on the same network). "
            "Step 9: Close the window. Done! The PC now has an IP address."
        ),
        "music": "Focused study music - 'Slow Motion' by Bensound",
        "animation": "Step-by-step numbered highlights on screenshot placeholder",
        "emotion": "Focused and practical - like a helpful tutor",
        "bullets": [
            "Step 1: Click on the PC",
            "Step 2: Click Desktop tab (top of window)",
            "Step 3: Click IP Configuration icon",
            "Step 4: Select Static (not DHCP)",
            "Step 5: Type IPv4 Address (e.g. 192.168.1.10)",
            "Step 6: Type Subnet Mask (255.255.255.0)",
            "Step 7: Type Default Gateway (router IP)",
            "Step 8: Close window - Done!"
        ],
        "bg": "FFFFFF"
    },
    {
        "num": 9,
        "title": "How to Place IP on a Router (Config Tab > Interface)",
        "script": (
            "Routers are a little different from PCs. Here is how to set an IP address on a "
            "router using only the GUI, no CLI. Step 1: Click on the Router in your workspace. "
            "A window pops up. Step 2: Click on the Config tab at the top. Step 3: On the left "
            "side, you will see a list of interfaces like GigabitEthernet0/0, GigabitEthernet0/1, "
            "etc. Click on the interface you want to configure (the one connected to your cable). "
            "Step 4: In the IP Configuration section on the right, type the IP Address for this "
            "interface. Step 5: Type the Subnet Mask. Step 6: VERY IMPORTANT - check the Port "
            "Status checkbox to turn the interface ON. The checkbox should say On. If it is off, "
            "the port is shut down and nothing will work. Step 7: Close the window. That interface "
            "now has an IP address and is active!"
        ),
        "music": "Tutorial progression music - 'Next Step' style",
        "animation": "Router interface config screen placeholder with arrows",
        "emotion": "Careful and detailed - each step is critical",
        "bullets": [
            "Step 1: Click on the Router",
            "Step 2: Click Config tab (top of window)",
            "Step 3: Click interface on left (e.g. Gig0/0)",
            "Step 4: Type IP Address (e.g. 192.168.1.1)",
            "Step 5: Type Subnet Mask (255.255.255.0)",
            "Step 6: CHECK Port Status = ON!",
            "Step 7: Close window - Interface is active!"
        ],
        "bg": "D6EAF8"
    },

    {
        "num": 10,
        "title": "How to Place IP on a Switch (Config Tab > VLAN1)",
        "script": (
            "Switches do not usually need an IP address for forwarding data. But if you want to "
            "manage the switch remotely or test connectivity TO the switch itself, you can give "
            "it an IP on VLAN1. Step 1: Click on the Switch. Step 2: Click the Config tab. "
            "Step 3: On the left side, under INTERFACE, look for VLAN1 (or click Settings then "
            "look for Interface VLAN1). Step 4: Type the IP Address. Step 5: Type the Subnet Mask. "
            "Step 6: Make sure the port status is ON. Step 7: Also set the Default Gateway under "
            "Settings on the left panel (this is the router IP so the switch can reach other networks). "
            "Step 8: Close the window."
        ),
        "music": "Focused study continues - 'Concentration' style",
        "animation": "Switch config screen placeholder with highlights",
        "emotion": "Thorough and precise - covering all bases",
        "bullets": [
            "Switches usually forward data without IP",
            "But VLAN1 IP allows management access",
            "Step 1: Click Switch > Config tab",
            "Step 2: Click VLAN1 on left panel",
            "Step 3: Type IP Address + Subnet Mask",
            "Step 4: Port Status = ON",
            "Step 5: Set Default Gateway under Settings"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 11,
        "title": "How to Use the Add Simple PDU (Mailbox Icon) to Test",
        "script": (
            "Now the most important testing tool! Forget the Command Prompt. Forget typing ping. "
            "We will use the Add Simple PDU tool instead. Look at the RIGHT side of the Packet "
            "Tracer workspace. There is a vertical toolbar with several icons. Look for the icon "
            "that looks like a CLOSED ENVELOPE (some people call it the mailbox icon). It is "
            "usually the first or second icon in that right toolbar. Step 1: Click on that "
            "envelope icon. Your cursor changes to a crosshair with an envelope. Step 2: Click on "
            "the SOURCE device (the device you want to send FROM). Step 3: Click on the "
            "DESTINATION device (the device you want to send TO). Step 4: Now look at the "
            "bottom-right panel. It shows a PDU list with columns: Last Status, Source, Destination, "
            "Type. Step 5: If it says Successful under Last Status, your devices can communicate! "
            "If it says Failed, something is wrong with your IP addresses, cables, or port status."
        ),
        "music": "Achievement music - 'Discovery' style",
        "animation": "Highlight right toolbar, then show click sequence",
        "emotion": "Empowering - this is the key skill!",
        "bullets": [
            "Right toolbar > Closed Envelope icon",
            "Step 1: Click the envelope/mailbox icon",
            "Step 2: Click SOURCE device (send from)",
            "Step 3: Click DESTINATION device (send to)",
            "Step 4: Check bottom-right PDU panel",
            "Successful = devices communicate!",
            "Failed = check IPs, cables, port status"
        ],
        "bg": "FDEBD0"
    },

    {
        "num": 12,
        "title": "Understanding the PDU Results Panel",
        "script": (
            "After you use the Add Simple PDU tool, let me explain what you see at the bottom. "
            "There is a panel that shows your PDU test results. Each row has: Last Status (Successful "
            "or Failed), Source (the device you clicked first), Destination (the device you clicked "
            "second), and Type (usually ICMP which is the ping protocol). You can also click "
            "Delete to remove old tests, or click Reset Simulation to start fresh. If you see "
            "Failed, do not panic! Just check: Are the IPs in the same network? Is the subnet mask "
            "correct? Is the router port turned ON? Is the default gateway set on the PCs? Fix these "
            "and try the PDU test again."
        ),
        "music": "Understanding music - 'Clarity' style",
        "animation": "Bottom panel zoom with annotations",
        "emotion": "Analytical and calm - interpreting results",
        "bullets": [
            "Bottom panel shows: Last Status column",
            "Successful = green checkmark = working!",
            "Failed = red X = something wrong",
            "Source and Destination columns shown",
            "Type = ICMP (ping protocol)",
            "Delete old tests or Reset Simulation",
            "Failed? Check IPs, subnet, gateway, port status"
        ],
        "bg": "FFFFFF"
    },
    {
        "num": 13,
        "title": "How to Download Cisco Packet Tracer",
        "script": (
            "Now let us learn how to get Packet Tracer on your computer. First, go to netacad.com. "
            "Create a free account with your name and email. After you sign up, look for the "
            "Resources section. Find Packet Tracer and click download. Choose the version for your "
            "computer: Windows, Mac, or Linux. It is about 200-300 MB. Make sure you have good internet!"
        ),
        "music": "Step-by-step tutorial music - 'Acoustic Breeze' by Bensound",
        "animation": "Numbered steps appear one by one from top",
        "emotion": "Helpful and patient - guiding step by step",
        "bullets": [
            "Go to netacad.com",
            "Create a FREE account",
            "Enter name and email",
            "Find Resources section",
            "Click Download Packet Tracer",
            "Choose: Windows / Mac / Linux",
            "File size: 200-300 MB"
        ],
        "bg": "D6EAF8"
    },

    {
        "num": 14,
        "title": "Installing and Opening Packet Tracer",
        "script": (
            "After download, double-click the file. On Windows: Next, Next, Finish. On Mac: drag "
            "to Applications. When you open it, log in with your NetAcad account. Takes 5-10 minutes. "
            "You are ready to build networks!"
        ),
        "music": "Calm progress music - 'Tomorrow' by Bensound",
        "animation": "Progress bar animation placeholder",
        "emotion": "Reassuring and simple - anyone can do this",
        "bullets": [
            "Double-click downloaded file",
            "Windows: Next, Next, Finish",
            "Mac: Drag to Applications",
            "Log in with NetAcad account",
            "Takes 5-10 minutes total",
            "You are ready to go!"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 15,
        "title": "The Workspace Overview",
        "script": (
            "When you open Packet Tracer, here is what you see. Menu bar at top. Toolbar below it. "
            "Big white workspace in the middle where you build. Device categories at the bottom. "
            "And on the RIGHT SIDE, the important vertical toolbar with the Add Simple PDU icon "
            "(the envelope/mailbox). Remember that right toolbar - we will use it a lot today!"
        ),
        "music": "Exploration music - 'Clear Day' by Bensound",
        "animation": "Pan across interface with right toolbar highlighted",
        "emotion": "Patient and guiding - like a tour guide",
        "bullets": [
            "Menu bar at the top",
            "Toolbar with quick buttons below",
            "Big white workspace in the middle",
            "Device categories at the bottom",
            "RIGHT SIDE: vertical toolbar",
            "Envelope icon = Add Simple PDU (our testing tool!)"
        ],
        "bg": "FDEBD0"
    },
    {
        "num": 16,
        "title": "Cable Types Quick Reference",
        "script": (
            "Before we start building scenarios, remember these cable rules. Straight-through cable: "
            "use for DIFFERENT devices (PC to Switch, Switch to Router). Cross-over cable: use for "
            "SIMILAR devices (PC to PC, Switch to Switch, Router to Router). If you are confused, "
            "just use the automatic cable tool (the lightning bolt with question mark). It picks the "
            "right cable for you!"
        ),
        "music": "Quick reference music - 'Organized' style",
        "animation": "Cable type comparison chart",
        "emotion": "Practical and helpful - quick reference",
        "bullets": [
            "Straight-through: PC to Switch",
            "Straight-through: Switch to Router",
            "Cross-over: PC to PC",
            "Cross-over: Switch to Switch",
            "Confused? Use automatic cable tool!",
            "Green dots = connection working"
        ],
        "bg": "FFFFFF"
    },

    # ===== SCENARIO 1: School Computer Lab in Manila (Slides 17-22) =====
    {
        "num": 17,
        "title": "Scenario 1: School Computer Lab in Manila - Overview",
        "script": (
            "Let us build our first real scenario! A public school in Manila has a computer lab. "
            "They have 4 PCs for students (we use 4 to keep it simple, but you can add more). "
            "They need one switch to connect all PCs, and one router for internet. Here is the "
            "IP plan: Router GigabitEthernet0/0 = 192.168.1.1. PC1 = 192.168.1.10. PC2 = 192.168.1.11. "
            "PC3 = 192.168.1.12. PC4 = 192.168.1.13. All use subnet mask 255.255.255.0. "
            "Default gateway for all PCs = 192.168.1.1 (the router)."
        ),
        "music": "Filipino-inspired gentle music - 'Maligaya' style acoustic",
        "animation": "Network diagram building up: Router > Switch > 4 PCs",
        "emotion": "Proud and relatable - showing Filipino context",
        "bullets": [
            "Devices: 1 Router + 1 Switch + 4 PCs",
            "Router Gig0/0: 192.168.1.1 / 255.255.255.0",
            "PC1: 192.168.1.10 / 255.255.255.0",
            "PC2: 192.168.1.11 / 255.255.255.0",
            "PC3: 192.168.1.12 / 255.255.255.0",
            "PC4: 192.168.1.13 / 255.255.255.0",
            "Default Gateway (all PCs): 192.168.1.1"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 18,
        "title": "Scenario 1: Manila School - Place Devices & Cables",
        "script": (
            "Step 1: From the bottom panel, click Network Devices > Routers. Drag a Router "
            "(like 1941 or 2911) onto the workspace near the top. Step 2: Click Network Devices > "
            "Switches. Drag a Switch (like 2960) to the middle area. Step 3: Click End Devices. "
            "Drag 4 PCs below the switch. Label them PC1, PC2, PC3, PC4. Step 4: Connect the "
            "Router to the Switch using a Straight-Through cable. Click Router > GigabitEthernet0/0, "
            "then click Switch > FastEthernet0/1. Step 5: Connect each PC to the Switch using "
            "Straight-Through cables. PC1 to Fa0/2, PC2 to Fa0/3, PC3 to Fa0/4, PC4 to Fa0/5. "
            "Wait for dots to turn green."
        ),
        "music": "Building music - 'Construction' style",
        "animation": "Devices appearing one by one, cables drawn",
        "emotion": "Active and building - creating something!",
        "bullets": [
            "Drag Router (1941) to top of workspace",
            "Drag Switch (2960) to middle",
            "Drag 4 PCs below the switch",
            "Cable: Router Gig0/0 to Switch Fa0/1 (Straight)",
            "Cable: PC1 to Switch Fa0/2 (Straight)",
            "Cable: PC2-Fa0/3, PC3-Fa0/4, PC4-Fa0/5",
            "Wait for green dots!"
        ],
        "bg": "D6EAF8"
    },

    {
        "num": 19,
        "title": "Scenario 1: Manila School - Set IP on Router (GUI Only)",
        "script": (
            "Now we set the IP on the router using ONLY the GUI. No CLI! Step 1: Click on the "
            "Router. The router window opens. Step 2: Click the Config tab at the top. Step 3: "
            "On the left panel, under INTERFACE, click GigabitEthernet0/0. Step 4: On the right "
            "side, you see IP Configuration section. In the IP Address field, type: 192.168.1.1. "
            "Step 5: In the Subnet Mask field, type: 255.255.255.0. Step 6: VERY IMPORTANT - "
            "look for the Port Status checkbox. Make sure it says ON (check the box). If it is "
            "off, the port is shut down and nothing works! Step 7: Close the router window. "
            "The router interface is now configured!"
        ),
        "music": "Detailed tutorial music - 'Focus' style",
        "animation": "Router config screen with fields highlighted step by step",
        "emotion": "Meticulous and careful - every detail matters",
        "bullets": [
            "Click Router > Config tab",
            "Left panel: click GigabitEthernet0/0",
            "IP Address: 192.168.1.1",
            "Subnet Mask: 255.255.255.0",
            "Port Status: CHECK the ON box!",
            "If Port Status is OFF = nothing works!",
            "Close window - Router is configured!"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 20,
        "title": "Scenario 1: Manila School - Set IP on All PCs (GUI Only)",
        "script": (
            "Now we set IP addresses on all 4 PCs. For PC1: Step 1: Click on PC1. Step 2: Click "
            "the Desktop tab. Step 3: Click IP Configuration. Step 4: Select Static (not DHCP). "
            "Step 5: IPv4 Address = 192.168.1.10. Step 6: Subnet Mask = 255.255.255.0. "
            "Step 7: Default Gateway = 192.168.1.1 (this is the router IP). Close the window. "
            "Repeat for PC2: same steps but IPv4 Address = 192.168.1.11. "
            "PC3: IPv4 Address = 192.168.1.12. PC4: IPv4 Address = 192.168.1.13. "
            "All PCs use the same Subnet Mask (255.255.255.0) and same Default Gateway (192.168.1.1)."
        ),
        "music": "Repetition practice music - 'Practice Makes Perfect' style",
        "animation": "PC config screen shown for each PC with values",
        "emotion": "Patient repetition - drilling the process",
        "bullets": [
            "PC1: Desktop > IP Config > Static",
            "  IP: 192.168.1.10 | Mask: 255.255.255.0 | GW: 192.168.1.1",
            "PC2: Desktop > IP Config > Static",
            "  IP: 192.168.1.11 | Mask: 255.255.255.0 | GW: 192.168.1.1",
            "PC3: IP: 192.168.1.12 | same mask & gateway",
            "PC4: IP: 192.168.1.13 | same mask & gateway",
            "All PCs: Gateway = 192.168.1.1 (Router IP)"
        ],
        "bg": "FDEBD0"
    },

    {
        "num": 21,
        "title": "Scenario 1: Manila School - Test with Add Simple PDU (Mailbox)",
        "script": (
            "Time to test! Remember, we do NOT open any terminal. We use the mailbox icon. "
            "Step 1: Look at the RIGHT side toolbar of the Packet Tracer workspace. Find the "
            "closed envelope icon (Add Simple PDU). Click it. Step 2: Your cursor changes to a "
            "crosshair with a small envelope. Step 3: Click on PC1 (the source). Step 4: Click "
            "on PC2 (the destination). Step 5: Look at the bottom-right panel. You should see "
            "a row appear: Last Status = Successful, Source = PC1, Destination = PC2, Type = ICMP. "
            "Successful means PC1 can talk to PC2! Step 6: Try again - click Add Simple PDU icon, "
            "click PC3, then click Router. If you see Successful, PC3 can reach the router too! "
            "Test all combinations to make sure your network works."
        ),
        "music": "Testing and celebration music - 'Moment of Truth' style",
        "animation": "Click sequence animation: envelope icon > PC1 > PC2 > result shown",
        "emotion": "Exciting - the moment of truth!",
        "bullets": [
            "RIGHT toolbar > Click envelope icon (Add Simple PDU)",
            "Cursor changes to crosshair + envelope",
            "Click PC1 (source), then Click PC2 (destination)",
            "Bottom panel: Last Status = Successful!",
            "Try PC3 to Router - should also be Successful",
            "Test all pairs: PC1>PC4, PC2>PC3, etc.",
            "All Successful = network is working perfectly!"
        ],
        "bg": "FFFFFF"
    },
    {
        "num": 22,
        "title": "Scenario 1: Manila School - Troubleshooting If Failed",
        "script": (
            "What if the PDU test shows Failed? Do not worry! Here is how to fix it. "
            "Check 1: Did you set the router port status to ON? Click Router > Config > Gig0/0 > "
            "make sure the On checkbox is checked. Check 2: Are all PCs set to Static (not DHCP)? "
            "Click each PC > Desktop > IP Configuration > make sure Static is selected. "
            "Check 3: Is the default gateway correct on every PC? It should be 192.168.1.1. "
            "Check 4: Are all IP addresses in the same network (192.168.1.x)? "
            "Check 5: Are the cables showing green dots? Red dots mean bad connection. "
            "Fix these issues and try the Add Simple PDU test again!"
        ),
        "music": "Problem-solving music - 'Detective' style",
        "animation": "Checklist appearing with green checkmarks",
        "emotion": "Supportive troubleshooter - we will fix this together",
        "bullets": [
            "Failed? Check these things:",
            "1. Router port status = ON?",
            "2. PCs set to Static (not DHCP)?",
            "3. Default Gateway = 192.168.1.1 on all PCs?",
            "4. All IPs in same network (192.168.1.x)?",
            "5. Cable dots are green (not red)?",
            "Fix and re-test with Add Simple PDU!"
        ],
        "bg": "D6EAF8"
    },

    # ===== SCENARIO 2: Barangay Health Center (Slides 23-28) =====
    {
        "num": 23,
        "title": "Scenario 2: Barangay Health Center - Overview & IP Plan",
        "script": (
            "Our second scenario! A barangay health center needs a network. They have 3 computers "
            "for health workers and 1 printer. They need a switch and a router for internet. "
            "IP Plan: Router Gig0/0 = 10.0.0.1. PC-Doctor = 10.0.0.10. PC-Nurse = 10.0.0.11. "
            "PC-Records = 10.0.0.12. Printer = 10.0.0.20. Subnet mask for all: 255.255.255.0. "
            "Default gateway for all devices: 10.0.0.1 (the router). Notice we are using a "
            "different network (10.0.0.x) from the school example (192.168.1.x). Both are valid!"
        ),
        "music": "Hopeful community music - 'Pagsisikap' style",
        "animation": "Community health center with network overlay",
        "emotion": "Caring and community-focused - showing real impact",
        "bullets": [
            "Devices: 1 Router + 1 Switch + 3 PCs + 1 Printer",
            "Network: 10.0.0.0/24",
            "Router Gig0/0: 10.0.0.1 / 255.255.255.0",
            "PC-Doctor: 10.0.0.10 / 255.255.255.0",
            "PC-Nurse: 10.0.0.11 / 255.255.255.0",
            "PC-Records: 10.0.0.12 / 255.255.255.0",
            "Printer: 10.0.0.20 / 255.255.255.0",
            "Gateway for all: 10.0.0.1"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 24,
        "title": "Scenario 2: Barangay Health - Place Devices & Cables",
        "script": (
            "Let us build it! Step 1: Drag a Router (2911) to the top. Step 2: Drag a Switch "
            "(2960) to the middle. Step 3: Drag 3 PCs to the bottom-left. Rename them: "
            "Right-click each PC > choose Display Name > type PC-Doctor, PC-Nurse, PC-Records. "
            "Step 4: Drag 1 Printer from End Devices to the bottom-right. Step 5: Connect Router "
            "Gig0/0 to Switch Fa0/1 using Straight-Through cable. Step 6: Connect PC-Doctor to "
            "Switch Fa0/2, PC-Nurse to Fa0/3, PC-Records to Fa0/4, and Printer to Fa0/5. "
            "All cables are Straight-Through because they go from end device to switch."
        ),
        "music": "Building music - 'Assembly' style",
        "animation": "Devices and cables appearing on workspace",
        "emotion": "Organized and methodical - building piece by piece",
        "bullets": [
            "Drag Router (2911) to top",
            "Drag Switch (2960) to middle",
            "Drag 3 PCs + 1 Printer to bottom",
            "Rename PCs: PC-Doctor, PC-Nurse, PC-Records",
            "Router Gig0/0 > Switch Fa0/1 (Straight cable)",
            "PCs to Switch Fa0/2-4, Printer to Fa0/5",
            "All Straight-Through cables (different devices)"
        ],
        "bg": "FDEBD0"
    },

    {
        "num": 25,
        "title": "Scenario 2: Barangay Health - Configure Router IP (GUI)",
        "script": (
            "Configure the router. Step 1: Click on the Router. Step 2: Click Config tab. "
            "Step 3: Click GigabitEthernet0/0 on the left panel. Step 4: Type IP Address: "
            "10.0.0.1. Step 5: Type Subnet Mask: 255.255.255.0. Step 6: Check the Port Status "
            "box to ON. Remember, if this checkbox is not checked, the interface is shut down "
            "and no data can pass through! This is the number one reason students get Failed "
            "results. Always turn the port ON. Step 7: Close the router window."
        ),
        "music": "Configuration music - 'Setup' style",
        "animation": "Router config highlighted fields",
        "emotion": "Careful and emphasizing - port status is key",
        "bullets": [
            "Click Router > Config tab",
            "Click GigabitEthernet0/0 (left panel)",
            "IP Address: 10.0.0.1",
            "Subnet Mask: 255.255.255.0",
            "Port Status: ON (CHECK THIS!)",
            "Port OFF = #1 reason for Failed tests",
            "Close window when done"
        ],
        "bg": "FFFFFF"
    },
    {
        "num": 26,
        "title": "Scenario 2: Barangay Health - Configure All PCs & Printer IP",
        "script": (
            "Now configure each device. PC-Doctor: Click > Desktop > IP Configuration > Static. "
            "IPv4 Address: 10.0.0.10. Subnet Mask: 255.255.255.0. Default Gateway: 10.0.0.1. Close. "
            "PC-Nurse: Same steps. IPv4 Address: 10.0.0.11. Same mask and gateway. Close. "
            "PC-Records: IPv4 Address: 10.0.0.12. Same mask and gateway. Close. "
            "Printer: Click on Printer > go to Config tab (printers use Config, not Desktop). "
            "Click on FastEthernet0 on the left. Type IP Address: 10.0.0.20. Subnet Mask: "
            "255.255.255.0. Default Gateway: 10.0.0.1. Make sure port status is ON. Close."
        ),
        "music": "Repetitive practice music - 'Drill' style",
        "animation": "Each device config shown in sequence",
        "emotion": "Thorough - every device must be configured",
        "bullets": [
            "PC-Doctor: Desktop > IP Config > Static",
            "  IP: 10.0.0.10 | Mask: 255.255.255.0 | GW: 10.0.0.1",
            "PC-Nurse: IP: 10.0.0.11 | same mask & GW",
            "PC-Records: IP: 10.0.0.12 | same mask & GW",
            "Printer: Config tab > FastEthernet0",
            "  IP: 10.0.0.20 | Mask: 255.255.255.0 | GW: 10.0.0.1",
            "Printer port status must be ON!"
        ],
        "bg": "D6EAF8"
    },

    {
        "num": 27,
        "title": "Scenario 2: Barangay Health - Test with Mailbox Icon",
        "script": (
            "Test time! Step 1: Click the Add Simple PDU icon (closed envelope) on the RIGHT "
            "toolbar. Step 2: Click on PC-Doctor (source). Step 3: Click on Printer (destination). "
            "Step 4: Check the bottom panel. It should show Successful! This means the doctor "
            "can print patient records. Step 5: Click the envelope icon again. Click PC-Nurse, "
            "then click PC-Records. Check bottom panel: Successful! The nurse can share files "
            "with the records computer. Step 6: Click envelope icon. Click PC-Records, then "
            "click Router. Successful means the records PC can reach the internet gateway!"
        ),
        "music": "Testing success music - 'Verification' style",
        "animation": "Envelope clicking sequence with success indicators",
        "emotion": "Satisfying - seeing everything work!",
        "bullets": [
            "Click envelope icon (right toolbar)",
            "Test 1: PC-Doctor > Printer = Successful!",
            "Test 2: PC-Nurse > PC-Records = Successful!",
            "Test 3: PC-Records > Router = Successful!",
            "All tests pass = network works perfectly!",
            "Health workers can print and share data!",
            "Technology serving the community!"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 28,
        "title": "Scenario 2: Barangay Health - Common Mistakes Here",
        "script": (
            "Common mistakes for this scenario: Mistake 1: Forgetting the Printer uses Config "
            "tab, not Desktop tab. PCs have Desktop tab for IP. Printers and Servers use Config "
            "tab. Mistake 2: Forgetting to set Default Gateway on the Printer. Without it, the "
            "printer can only talk to devices on the same switch, not through the router. "
            "Mistake 3: Using DHCP instead of Static on PCs. If you leave it on DHCP but there "
            "is no DHCP server, the PC gets no IP address! Always use Static for our exercises. "
            "Mistake 4: Wrong subnet mask. If one device has 255.255.0.0 and another has "
            "255.255.255.0, they might not communicate properly."
        ),
        "music": "Warning and tips music - 'Watch Out' style",
        "animation": "Common mistakes with X marks and corrections",
        "emotion": "Preventive - catching mistakes before they happen",
        "bullets": [
            "Printer uses Config tab (not Desktop!)",
            "ALWAYS set Default Gateway on printers too",
            "Use Static, not DHCP (no DHCP server here)",
            "Subnet mask must match on ALL devices",
            "255.255.255.0 on everything in same network",
            "Mismatch = communication failure!",
            "Double-check every device before testing"
        ],
        "bg": "FDEBD0"
    },

    # ===== SCENARIO 3: Sari-Sari Store POS Network (Slides 29-33) =====
    {
        "num": 29,
        "title": "Scenario 3: Sari-Sari Store POS Network - Overview & IP Plan",
        "script": (
            "Scenario 3! A sari-sari store wants a point-of-sale system and WiFi for customers. "
            "Devices: 1 Router, 1 Switch, 1 PC (cash register), 1 Laptop (owner), 1 Wireless "
            "Access Point, 1 Smartphone (customer WiFi). IP Plan: Router Gig0/0 = 172.16.0.1. "
            "PC-CashRegister = 172.16.0.10. Laptop-Owner = 172.16.0.11. Access Point does not "
            "need an IP for basic operation. Smartphone-Customer = 172.16.0.50. "
            "All subnet masks: 255.255.255.0. All gateways: 172.16.0.1. This time we use the "
            "172.16.0.x network to show you another valid IP range!"
        ),
        "music": "Cheerful entrepreneurial music - 'Negosyo Beat'",
        "animation": "Store diagram with network overlay",
        "emotion": "Entrepreneurial and encouraging - small business tech",
        "bullets": [
            "Devices: Router + Switch + PC + Laptop + AP + Phone",
            "Network: 172.16.0.0/24",
            "Router Gig0/0: 172.16.0.1",
            "PC-CashRegister: 172.16.0.10",
            "Laptop-Owner: 172.16.0.11",
            "Smartphone-Customer: 172.16.0.50",
            "All masks: 255.255.255.0 | GW: 172.16.0.1"
        ],
        "bg": "FFFFFF"
    },
    {
        "num": 30,
        "title": "Scenario 3: Sari-Sari Store - Build & Cable the Network",
        "script": (
            "Build it! Step 1: Drag Router (1941) to top. Step 2: Drag Switch (2960) to middle. "
            "Step 3: Drag a PC to bottom-left, rename to PC-CashRegister. Step 4: Drag a Laptop "
            "to bottom-middle, rename to Laptop-Owner. Step 5: From Network Devices > Wireless, "
            "drag an Access Point (like AP-PT) to the right side. Step 6: From End Devices > "
            "Wireless, drag a Smartphone below the access point. Step 7: Connect Router Gig0/0 "
            "to Switch Fa0/1 (Straight). Connect PC to Switch Fa0/2 (Straight). Connect Laptop "
            "to Switch Fa0/3 (Straight). Connect Access Point Port1 to Switch Fa0/4 (Straight). "
            "The Smartphone connects wirelessly to the Access Point automatically!"
        ),
        "music": "Building and connecting music - 'Wire It Up' style",
        "animation": "Wired and wireless connections shown",
        "emotion": "Fun and practical - building a real business network",
        "bullets": [
            "Router (1941) at top, Switch (2960) middle",
            "PC-CashRegister + Laptop-Owner at bottom",
            "Access Point (AP-PT) on right side",
            "Smartphone below AP (connects wirelessly)",
            "Router Gig0/0 > Switch Fa0/1 (Straight)",
            "PC>Fa0/2, Laptop>Fa0/3, AP>Fa0/4 (all Straight)",
            "Phone connects to AP wirelessly (automatic!)"
        ],
        "bg": "D6EAF8"
    },

    {
        "num": 31,
        "title": "Scenario 3: Sari-Sari Store - Configure All IPs (GUI Only)",
        "script": (
            "Configure everything using GUI only! Router: Click Router > Config > GigabitEthernet0/0. "
            "IP: 172.16.0.1, Mask: 255.255.255.0, Port Status ON. Close. "
            "PC-CashRegister: Click > Desktop > IP Configuration > Static. IP: 172.16.0.10, "
            "Mask: 255.255.255.0, Gateway: 172.16.0.1. Close. "
            "Laptop-Owner: Click > Desktop > IP Configuration > Static. IP: 172.16.0.11, "
            "Mask: 255.255.255.0, Gateway: 172.16.0.1. Close. "
            "Smartphone-Customer: Click on the Smartphone > go to Config tab > Wireless0 "
            "interface on the left. Set IP: 172.16.0.50, Mask: 255.255.255.0, Gateway: 172.16.0.1. "
            "The Access Point itself does not need an IP for basic bridging."
        ),
        "music": "Configuration sequence music - 'One by One' style",
        "animation": "Each device config screen shown in sequence",
        "emotion": "Systematic - configure one device at a time",
        "bullets": [
            "Router: Config > Gig0/0 > IP:172.16.0.1 > ON",
            "PC-CashRegister: Desktop > IP Config > Static",
            "  IP:172.16.0.10 | Mask:255.255.255.0 | GW:172.16.0.1",
            "Laptop-Owner: Desktop > IP Config > Static",
            "  IP:172.16.0.11 | Mask:255.255.255.0 | GW:172.16.0.1",
            "Smartphone: Config > Wireless0",
            "  IP:172.16.0.50 | Mask:255.255.255.0 | GW:172.16.0.1"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 32,
        "title": "Scenario 3: Sari-Sari Store - Test with Mailbox Icon",
        "script": (
            "Let us test! Step 1: Click the Add Simple PDU (envelope) icon on the right toolbar. "
            "Step 2: Click PC-CashRegister, then click Laptop-Owner. Check bottom panel: "
            "Successful! The cash register can communicate with the owner laptop. "
            "Step 3: Click envelope icon again. Click Laptop-Owner, then click Smartphone-Customer. "
            "Successful! Wired and wireless devices can communicate. "
            "Step 4: Click envelope icon. Click Smartphone-Customer, then click Router. "
            "Successful! Even the wireless customer can reach the internet gateway. "
            "This proves your entire sari-sari store network is working!"
        ),
        "music": "Victory testing music - 'All Connected' style",
        "animation": "PDU envelope traveling between devices",
        "emotion": "Triumphant - wireless and wired both work!",
        "bullets": [
            "Click envelope icon (right toolbar)",
            "Test 1: PC-CashRegister > Laptop-Owner = Successful!",
            "Test 2: Laptop-Owner > Smartphone = Successful!",
            "Test 3: Smartphone > Router = Successful!",
            "Wired AND wireless both communicate!",
            "Entire store network is operational!",
            "Small business technology success!"
        ],
        "bg": "FDEBD0"
    },
    {
        "num": 33,
        "title": "Scenario 3: Sari-Sari Store - Wireless Tips",
        "script": (
            "Important tips for the wireless part: Tip 1: The Smartphone must have its wireless "
            "interface turned ON in Config > Wireless0. Tip 2: The Access Point SSID (network name) "
            "and the Smartphone must use the same SSID. By default they both use 'Default' so "
            "it usually works automatically. Tip 3: If the phone is not connecting to the AP, "
            "click the Smartphone > Config > Wireless0, and make sure the SSID field matches "
            "the AP SSID. Tip 4: The Access Point acts as a bridge. It does not need its own IP. "
            "It just passes traffic between wireless devices and the wired switch."
        ),
        "music": "Tips and hints music - 'Wireless Waves' style",
        "animation": "Wireless connection animation with SSID labels",
        "emotion": "Helpful extra details - making wireless work",
        "bullets": [
            "Smartphone wireless must be ON",
            "SSID must match between Phone and AP",
            "Default SSID = 'Default' (works automatically)",
            "If not connecting: check SSID match",
            "Access Point = bridge (no IP needed)",
            "AP passes traffic: wireless to wired",
            "All devices share same IP network"
        ],
        "bg": "FFFFFF"
    },

    # ===== SCENARIO 4: University Campus in Cebu - TWO networks (Slides 34-39) =====
    {
        "num": 34,
        "title": "Scenario 4: University Campus in Cebu - Overview & IP Plan",
        "script": (
            "Now a bigger challenge! A university in Cebu has TWO buildings: Library and Science. "
            "Each building has its own network connected by a router. This teaches you how a "
            "router connects TWO different networks. IP Plan: Router Gig0/0 (Library side) = "
            "192.168.10.1. Router Gig0/1 (Science side) = 192.168.20.1. Library-PC1 = 192.168.10.10. "
            "Library-PC2 = 192.168.10.11. Science-PC1 = 192.168.20.10. Science-PC2 = 192.168.20.11. "
            "Library PCs gateway: 192.168.10.1. Science PCs gateway: 192.168.20.1. "
            "Two switches: one for each building."
        ),
        "music": "Academic achievement music - 'Graduation Day' style",
        "animation": "Two-building campus diagram",
        "emotion": "Ambitious - stepping up to a bigger network",
        "bullets": [
            "TWO buildings = TWO networks connected by 1 Router",
            "Router Gig0/0 (Library): 192.168.10.1",
            "Router Gig0/1 (Science): 192.168.20.1",
            "Library-PC1: 192.168.10.10 | GW: 192.168.10.1",
            "Library-PC2: 192.168.10.11 | GW: 192.168.10.1",
            "Science-PC1: 192.168.20.10 | GW: 192.168.20.1",
            "Science-PC2: 192.168.20.11 | GW: 192.168.20.1"
        ],
        "bg": "D6EAF8"
    },
    {
        "num": 35,
        "title": "Scenario 4: Cebu University - Build the Two-Network Topology",
        "script": (
            "Build it! Step 1: Drag 1 Router (2911) to the center-top. This router has TWO "
            "GigabitEthernet ports. Step 2: Drag Switch-Library (2960) to the left-middle. "
            "Step 3: Drag Switch-Science (2960) to the right-middle. Step 4: Drag Library-PC1 "
            "and Library-PC2 below the left switch. Step 5: Drag Science-PC1 and Science-PC2 "
            "below the right switch. Step 6: Connect Router Gig0/0 to Switch-Library Fa0/1 "
            "(Straight). Connect Router Gig0/1 to Switch-Science Fa0/1 (Straight). "
            "Step 7: Connect Library-PC1 to Switch-Library Fa0/2, Library-PC2 to Fa0/3. "
            "Step 8: Connect Science-PC1 to Switch-Science Fa0/2, Science-PC2 to Fa0/3."
        ),
        "music": "Building complex networks music - 'Level Up' style",
        "animation": "Two-sided network building up from center",
        "emotion": "Confident - we can handle bigger networks!",
        "bullets": [
            "1 Router (2911) center-top",
            "Switch-Library (left), Switch-Science (right)",
            "2 PCs under each switch",
            "Router Gig0/0 > Switch-Library Fa0/1",
            "Router Gig0/1 > Switch-Science Fa0/1",
            "Library PCs > Switch-Library Fa0/2-3",
            "Science PCs > Switch-Science Fa0/2-3"
        ],
        "bg": "D5F5E3"
    },

    {
        "num": 36,
        "title": "Scenario 4: Cebu University - Configure Router TWO Interfaces",
        "script": (
            "The router needs TWO interfaces configured because it connects TWO networks. "
            "Step 1: Click Router > Config tab. Step 2: Click GigabitEthernet0/0 on the left. "
            "IP Address: 192.168.10.1. Subnet Mask: 255.255.255.0. Port Status: ON. "
            "Step 3: Now click GigabitEthernet0/1 on the left. IP Address: 192.168.20.1. "
            "Subnet Mask: 255.255.255.0. Port Status: ON. Step 4: Close the router window. "
            "BOTH interfaces must be ON! If you forget to turn on Gig0/1, the Science building "
            "will not be able to reach the Library building through the router."
        ),
        "music": "Dual configuration music - 'Two Sides' style",
        "animation": "Router with two interfaces highlighted",
        "emotion": "Precise - both interfaces matter equally",
        "bullets": [
            "Router has TWO interfaces to configure!",
            "Gig0/0: 192.168.10.1 / 255.255.255.0 / ON",
            "Gig0/1: 192.168.20.1 / 255.255.255.0 / ON",
            "BOTH must have Port Status = ON!",
            "Gig0/0 connects to Library network",
            "Gig0/1 connects to Science network",
            "Router routes traffic between them"
        ],
        "bg": "FDEBD0"
    },
    {
        "num": 37,
        "title": "Scenario 4: Cebu University - Configure All PCs",
        "script": (
            "Now configure the 4 PCs. Library-PC1: Click > Desktop > IP Configuration > Static. "
            "IP: 192.168.10.10, Mask: 255.255.255.0, Gateway: 192.168.10.1. Close. "
            "Library-PC2: IP: 192.168.10.11, Mask: 255.255.255.0, Gateway: 192.168.10.1. Close. "
            "Science-PC1: IP: 192.168.20.10, Mask: 255.255.255.0, Gateway: 192.168.20.1. Close. "
            "Science-PC2: IP: 192.168.20.11, Mask: 255.255.255.0, Gateway: 192.168.20.1. Close. "
            "NOTICE: Library PCs use gateway 192.168.10.1, Science PCs use gateway 192.168.20.1. "
            "Each side uses the router interface on ITS OWN network as the gateway!"
        ),
        "music": "Four-device configuration music - 'Sequence' style",
        "animation": "Four PC configs shown with different gateway highlights",
        "emotion": "Clear distinction - different networks, different gateways",
        "bullets": [
            "Library-PC1: 192.168.10.10 | GW: 192.168.10.1",
            "Library-PC2: 192.168.10.11 | GW: 192.168.10.1",
            "Science-PC1: 192.168.20.10 | GW: 192.168.20.1",
            "Science-PC2: 192.168.20.11 | GW: 192.168.20.1",
            "Library GW = Router Gig0/0 IP",
            "Science GW = Router Gig0/1 IP",
            "Gateway = YOUR side of the router!"
        ],
        "bg": "FFFFFF"
    },

    {
        "num": 38,
        "title": "Scenario 4: Cebu University - Test ACROSS Networks with PDU",
        "script": (
            "Now the exciting part! Can Library-PC1 talk to Science-PC1? They are on DIFFERENT "
            "networks! Step 1: Click Add Simple PDU (envelope icon, right toolbar). Step 2: "
            "Click Library-PC1. Step 3: Click Science-PC1. Step 4: Check bottom panel. "
            "It should show Successful! The router is routing the message from network "
            "192.168.10.x to network 192.168.20.x. Step 5: Test Library-PC2 to Science-PC2. "
            "Successful! Step 6: Test Science-PC1 to Library-PC1. Also Successful! "
            "The router allows BOTH directions. This is how buildings on a campus communicate!"
        ),
        "music": "Cross-network success music - 'Bridge Built' style",
        "animation": "PDU traveling across router from one network to another",
        "emotion": "Amazed - crossing network boundaries!",
        "bullets": [
            "Click envelope icon (right toolbar)",
            "Test: Library-PC1 > Science-PC1 = Successful!",
            "Router ROUTES between 192.168.10.x and 192.168.20.x",
            "Test: Library-PC2 > Science-PC2 = Successful!",
            "Test: Science-PC1 > Library-PC1 = Successful!",
            "Both directions work through the router!",
            "Campus buildings can communicate!"
        ],
        "bg": "D6EAF8"
    },
    {
        "num": 39,
        "title": "Scenario 4: Cebu University - Why Gateways Matter",
        "script": (
            "Let me explain WHY the gateway is so important in this scenario. When Library-PC1 "
            "wants to talk to Science-PC1, it checks: Is 192.168.20.10 on my network (192.168.10.x)? "
            "No! So it sends the message to its Default Gateway (192.168.10.1 = the router). "
            "The router looks at the destination and says: 192.168.20.10 is on my Gig0/1 network! "
            "It forwards the message out Gig0/1 to the Science switch, which delivers it to "
            "Science-PC1. If you forget to set the gateway on Library-PC1, it does not know WHERE "
            "to send messages for other networks. The PDU test will show Failed!"
        ),
        "music": "Understanding music - 'Light Bulb Moment' style",
        "animation": "Packet journey visualization through router",
        "emotion": "Enlightening - now it all makes sense!",
        "bullets": [
            "PC checks: Is destination on MY network?",
            "If NO: send to Default Gateway (router)",
            "Router checks its interfaces for destination",
            "Router forwards to correct network",
            "No gateway = PC does not know where to send!",
            "No gateway = Failed PDU test!",
            "Gateway = door to other networks"
        ],
        "bg": "D5F5E3"
    },

    # ===== SCENARIO 5: Government Office in Davao (Slides 40-44) =====
    {
        "num": 40,
        "title": "Scenario 5: Government Office in Davao - Overview & IP Plan",
        "script": (
            "Scenario 5! A government office in Davao needs a secure network with TWO zones: "
            "Public (for citizen services) and Private (for employee records). We use one router "
            "with two interfaces to separate them. IP Plan: Router Gig0/0 (Public zone) = "
            "192.168.100.1. Router Gig0/1 (Private zone) = 192.168.200.1. Public-PC1 (citizen "
            "kiosk) = 192.168.100.10. Public-PC2 = 192.168.100.11. Private-PC1 (employee) = "
            "192.168.200.10. Private-PC2 (records) = 192.168.200.11. Private-Server = 192.168.200.100. "
            "Notice the networks are far apart: 192.168.100.x vs 192.168.200.x for clear separation!"
        ),
        "music": "Professional and secure - 'Official Business' style",
        "animation": "Government building split into public and private zones",
        "emotion": "Serious but accessible - security matters",
        "bullets": [
            "TWO zones: Public (citizens) + Private (employees)",
            "Router Gig0/0 (Public): 192.168.100.1",
            "Router Gig0/1 (Private): 192.168.200.1",
            "Public-PC1: 192.168.100.10 | GW: 192.168.100.1",
            "Public-PC2: 192.168.100.11 | GW: 192.168.100.1",
            "Private-PC1: 192.168.200.10 | GW: 192.168.200.1",
            "Private-PC2: 192.168.200.11 | GW: 192.168.200.1",
            "Private-Server: 192.168.200.100 | GW: 192.168.200.1"
        ],
        "bg": "FDEBD0"
    },
    {
        "num": 41,
        "title": "Scenario 5: Davao Gov Office - Build the Topology",
        "script": (
            "Build it! Step 1: Drag Router (2911) to the top-center. Step 2: Drag Switch-Public "
            "to the left-middle area. Step 3: Drag Switch-Private to the right-middle area. "
            "Step 4: Drag Public-PC1 and Public-PC2 below the left switch. Step 5: Drag "
            "Private-PC1, Private-PC2, and a Server below the right switch. The Server is found "
            "in End Devices > Servers. Rename it Private-Server. Step 6: Cable Router Gig0/0 to "
            "Switch-Public Fa0/1 (Straight). Router Gig0/1 to Switch-Private Fa0/1 (Straight). "
            "Step 7: Connect all PCs and Server to their respective switches using Straight-Through."
        ),
        "music": "Government building music - 'Structured Order' style",
        "animation": "Dual-zone network building up",
        "emotion": "Organized - clear separation between zones",
        "bullets": [
            "Router (2911) top-center",
            "Switch-Public (left), Switch-Private (right)",
            "2 PCs under Public switch",
            "2 PCs + 1 Server under Private switch",
            "Router Gig0/0 > Switch-Public Fa0/1",
            "Router Gig0/1 > Switch-Private Fa0/1",
            "All end devices to switches (Straight cable)"
        ],
        "bg": "FFFFFF"
    },

    {
        "num": 42,
        "title": "Scenario 5: Davao Gov Office - Configure Router & All Devices",
        "script": (
            "Configure the router first. Click Router > Config > Gig0/0: IP 192.168.100.1, "
            "Mask 255.255.255.0, Port ON. Click Gig0/1: IP 192.168.200.1, Mask 255.255.255.0, "
            "Port ON. Close. Now the PCs. Public-PC1: Desktop > IP Config > Static. "
            "IP: 192.168.100.10, Mask: 255.255.255.0, Gateway: 192.168.100.1. "
            "Public-PC2: IP: 192.168.100.11, same mask and gateway. "
            "Private-PC1: Desktop > IP Config > Static. IP: 192.168.200.10, Mask: 255.255.255.0, "
            "Gateway: 192.168.200.1. Private-PC2: IP: 192.168.200.11, same mask and gateway. "
            "Private-Server: Click Server > Config tab > FastEthernet0. IP: 192.168.200.100, "
            "Mask: 255.255.255.0, Gateway: 192.168.200.1, Port ON."
        ),
        "music": "Multi-device config music - 'All Hands' style",
        "animation": "Device configs shown in rapid sequence",
        "emotion": "Efficient - configuring many devices smoothly",
        "bullets": [
            "Router: Gig0/0=192.168.100.1 ON, Gig0/1=192.168.200.1 ON",
            "Public-PC1: 192.168.100.10 | GW: 192.168.100.1",
            "Public-PC2: 192.168.100.11 | GW: 192.168.100.1",
            "Private-PC1: 192.168.200.10 | GW: 192.168.200.1",
            "Private-PC2: 192.168.200.11 | GW: 192.168.200.1",
            "Server: Config>Fa0 > 192.168.200.100 | GW: 192.168.200.1",
            "Server uses Config tab (like Printer!)"
        ],
        "bg": "D6EAF8"
    },
    {
        "num": 43,
        "title": "Scenario 5: Davao Gov Office - Test All Connections with PDU",
        "script": (
            "Test everything with the mailbox icon! Click the Add Simple PDU envelope on the "
            "right toolbar. Test 1: Public-PC1 > Public-PC2. Successful! Same network. "
            "Test 2: Private-PC1 > Private-Server. Successful! Employee can access records. "
            "Test 3: Public-PC1 > Private-Server. Successful! Wait - is this a security problem? "
            "Yes! In real life you would add firewall rules. But for now, the router allows all "
            "traffic by default. This teaches you that basic routing connects everything. "
            "Security (blocking certain traffic) is an advanced topic for later."
        ),
        "music": "Testing and discovery music - 'What Happens' style",
        "animation": "PDU tests with security question mark",
        "emotion": "Curious - discovering security implications",
        "bullets": [
            "Test 1: Public-PC1 > Public-PC2 = Successful!",
            "Test 2: Private-PC1 > Private-Server = Successful!",
            "Test 3: Public-PC1 > Private-Server = Successful!",
            "Wait - citizens can reach private server?!",
            "Yes! Router allows ALL traffic by default",
            "Security requires firewall rules (advanced topic)",
            "Basic routing = everything connected"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 44,
        "title": "Scenario 5: Davao Gov Office - Security Awareness",
        "script": (
            "This scenario teaches us something important: just separating networks with different "
            "IP ranges does NOT automatically block traffic. The router happily forwards everything. "
            "In real government offices, they use Access Control Lists (ACLs) on the router to "
            "block unwanted traffic. They also use firewalls. But that is an advanced topic. "
            "For now, remember: a router CONNECTS networks. If you want to BLOCK traffic between "
            "networks, you need extra configuration. Our PDU tests proved that both zones CAN talk. "
            "Whether they SHOULD talk is a security decision."
        ),
        "music": "Wisdom and awareness music - 'Think About It' style",
        "animation": "Security lock with question mark",
        "emotion": "Thought-provoking - introducing security concepts",
        "bullets": [
            "Different IP ranges alone does NOT block traffic",
            "Router connects and forwards by default",
            "ACLs (Access Control Lists) can block traffic",
            "Firewalls add extra security",
            "These are advanced topics for later",
            "Router CONNECTS, ACLs/Firewalls BLOCK",
            "PDU tests help you verify what CAN communicate"
        ],
        "bg": "FDEBD0"
    },

    # ===== SCENARIO 6: Internet Cafe / Pisonet (Slides 45-50) =====
    {
        "num": 45,
        "title": "Scenario 6: Internet Cafe (Pisonet) - Overview & IP Plan",
        "script": (
            "Our final scenario! A pisonet in your barangay has 6 computers for customers plus "
            "1 server for managing time and payments. All connected to one switch and one router. "
            "IP Plan: Router Gig0/0 = 192.168.50.1. Server-Timer = 192.168.50.2. "
            "Customer-PC1 = 192.168.50.10. Customer-PC2 = 192.168.50.11. Customer-PC3 = 192.168.50.12. "
            "Customer-PC4 = 192.168.50.13. Customer-PC5 = 192.168.50.14. Customer-PC6 = 192.168.50.15. "
            "All subnet masks: 255.255.255.0. All gateways: 192.168.50.1."
        ),
        "music": "Gaming and fun music - 'Pixel Adventure' style",
        "animation": "Internet cafe layout with numbered PCs",
        "emotion": "Fun and relatable - connecting to everyday Filipino life",
        "bullets": [
            "Pisonet: 6 Customer PCs + 1 Server + Switch + Router",
            "Network: 192.168.50.0/24",
            "Router Gig0/0: 192.168.50.1",
            "Server-Timer: 192.168.50.2",
            "Customer-PC1: 192.168.50.10",
            "Customer-PC2 through PC6: .11 to .15",
            "All masks: 255.255.255.0 | GW: 192.168.50.1"
        ],
        "bg": "FFFFFF"
    },
    {
        "num": 46,
        "title": "Scenario 6: Pisonet - Build the Full Topology",
        "script": (
            "Build this network! Step 1: Drag Router (1941) to top. Step 2: Drag Switch (2960) "
            "to middle. Note: you need a switch with enough ports for 8 connections (6 PCs + "
            "1 server + 1 router). The 2960 has 24 FastEthernet ports, plenty! Step 3: Drag "
            "1 Server from End Devices, place near the switch, rename Server-Timer. Step 4: "
            "Drag 6 PCs below the switch. Rename them Customer-PC1 through Customer-PC6. "
            "Step 5: Connect Router Gig0/0 to Switch Fa0/1 (Straight). Server to Switch Fa0/2 "
            "(Straight). PC1 to Fa0/3, PC2 to Fa0/4, PC3 to Fa0/5, PC4 to Fa0/6, PC5 to Fa0/7, "
            "PC6 to Fa0/8. All Straight-Through cables."
        ),
        "music": "Assembly line music - 'Mass Production' style",
        "animation": "Many devices being placed in organized rows",
        "emotion": "Industrious - building a real business",
        "bullets": [
            "Router (1941) at top",
            "Switch (2960, 24 ports) in middle",
            "1 Server + 6 PCs below switch",
            "Rename: Server-Timer, Customer-PC1-6",
            "Router Gig0/0 > Switch Fa0/1",
            "Server > Fa0/2",
            "PC1>Fa0/3, PC2>Fa0/4 ... PC6>Fa0/8",
            "All Straight-Through cables"
        ],
        "bg": "D6EAF8"
    },

    {
        "num": 47,
        "title": "Scenario 6: Pisonet - Configure Router & Server",
        "script": (
            "Configure the router first. Click Router > Config > GigabitEthernet0/0. "
            "IP: 192.168.50.1, Mask: 255.255.255.0, Port Status ON. Close. "
            "Now the server. Click Server-Timer > Config tab > FastEthernet0 (left panel). "
            "IP Address: 192.168.50.2. Subnet Mask: 255.255.255.0. Default Gateway: 192.168.50.1. "
            "Port Status: ON. Important note about servers: Like printers, servers use the "
            "Config tab > interface for IP configuration, NOT the Desktop tab. The Desktop tab "
            "on servers is for running services (like DNS, DHCP, HTTP). For IP addressing, "
            "always go to Config tab for servers!"
        ),
        "music": "Server configuration music - 'System Online' style",
        "animation": "Router then Server config screens",
        "emotion": "Technical but clear - servers are special",
        "bullets": [
            "Router: Config > Gig0/0 > 192.168.50.1 > ON",
            "Server: Config tab > FastEthernet0",
            "Server IP: 192.168.50.2",
            "Server Mask: 255.255.255.0",
            "Server Gateway: 192.168.50.1",
            "Server Port Status: ON",
            "Servers use Config tab for IP (not Desktop!)"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 48,
        "title": "Scenario 6: Pisonet - Configure All 6 Customer PCs",
        "script": (
            "Now all 6 customer PCs. Same process for each! Click PC > Desktop tab > "
            "IP Configuration > Static. Customer-PC1: IP 192.168.50.10, Mask 255.255.255.0, "
            "Gateway 192.168.50.1. Customer-PC2: IP 192.168.50.11, same mask, same gateway. "
            "Customer-PC3: IP 192.168.50.12. Customer-PC4: IP 192.168.50.13. "
            "Customer-PC5: IP 192.168.50.14. Customer-PC6: IP 192.168.50.15. "
            "All use the SAME subnet mask (255.255.255.0) and SAME gateway (192.168.50.1). "
            "Only the last number of the IP changes for each PC! That is the pattern."
        ),
        "music": "Repetition and rhythm music - 'Pattern' style",
        "animation": "Six PC configs in rapid sequence showing pattern",
        "emotion": "Rhythmic and pattern-focused - see the system",
        "bullets": [
            "All PCs: Desktop > IP Config > Static",
            "PC1: 192.168.50.10 | PC2: 192.168.50.11",
            "PC3: 192.168.50.12 | PC4: 192.168.50.13",
            "PC5: 192.168.50.14 | PC6: 192.168.50.15",
            "All Masks: 255.255.255.0",
            "All Gateways: 192.168.50.1",
            "Only the LAST number changes per PC!"
        ],
        "bg": "FDEBD0"
    },

    {
        "num": 49,
        "title": "Scenario 6: Pisonet - Test All Connections with Mailbox",
        "script": (
            "Let us test the pisonet! Click the Add Simple PDU (envelope icon) on the right "
            "toolbar. Test 1: Customer-PC1 > Server-Timer. Successful! The customer can connect "
            "to the timer server. Test 2: Customer-PC3 > Customer-PC5. Successful! Customers "
            "can play LAN games with each other. Test 3: Customer-PC6 > Router. Successful! "
            "Customers can reach the internet. Test 4: Server-Timer > Customer-PC2. Successful! "
            "The server can push time limits to customer PCs. All tests pass! Your pisonet "
            "network is ready for business! Congratulations!"
        ),
        "music": "Business success music - 'Open for Business' style",
        "animation": "Multiple PDU tests shown passing",
        "emotion": "Celebratory - business is ready!",
        "bullets": [
            "Click envelope icon (right toolbar)",
            "Test 1: PC1 > Server-Timer = Successful!",
            "Test 2: PC3 > PC5 = Successful! (LAN gaming)",
            "Test 3: PC6 > Router = Successful! (internet)",
            "Test 4: Server > PC2 = Successful! (management)",
            "All tests pass!",
            "Pisonet is ready for business!"
        ],
        "bg": "FFFFFF"
    },
    {
        "num": 50,
        "title": "Scenario 6: Pisonet - IP Address Planning Tips",
        "script": (
            "Here are tips for planning IP addresses in any pisonet. Tip 1: Keep your numbering "
            "organized. Router = .1, Server = .2 to .9, Customer PCs = .10 onwards. This makes "
            "it easy to find devices. Tip 2: Write your IP plan on paper BEFORE you start "
            "configuring. Tip 3: The subnet mask 255.255.255.0 allows 254 devices on one network. "
            "That is more than enough for any pisonet. Tip 4: If you add more PCs later, just "
            "give them the next numbers: .16, .17, .18, etc. Tip 5: Never give two devices the "
            "same IP address! That causes a conflict and neither device will work properly."
        ),
        "music": "Planning and organization music - 'Strategy' style",
        "animation": "IP address table with organized numbering",
        "emotion": "Organized and strategic - plan before you build",
        "bullets": [
            "Router = .1, Servers = .2-.9, PCs = .10+",
            "Write IP plan on paper FIRST",
            "255.255.255.0 allows 254 devices",
            "Add more PCs: just use next numbers",
            "NEVER use same IP on two devices!",
            "Duplicate IP = conflict = both devices fail",
            "Good planning = easy troubleshooting"
        ],
        "bg": "D6EAF8"
    },

    # ===== COMMON PROBLEMS & TIPS (Slides 51-56) =====
    {
        "num": 51,
        "title": "Common Mistakes: Port Status Not Turned ON",
        "script": (
            "The NUMBER ONE mistake students make: forgetting to turn the router port ON! "
            "When you configure a router interface in the Config tab, there is a checkbox "
            "that says Port Status. If it is not checked (OFF), the interface is shut down. "
            "No data passes through. Your cables might show green dots because the physical "
            "connection works, but logically the port is disabled. ALWAYS check Port Status = ON "
            "after setting the IP address. This applies to EVERY interface on EVERY router. "
            "If your PDU test shows Failed, check this FIRST!"
        ),
        "music": "Urgent warning music - 'Do Not Forget' style",
        "animation": "Big flashing checkbox animation",
        "emotion": "Urgent - this is the most common mistake!",
        "bullets": [
            "#1 MISTAKE: Router port is OFF!",
            "Config > Interface > Port Status must be ON",
            "OFF = interface is shut down",
            "Green cable dots but still fails? = Port OFF!",
            "Check EVERY router interface",
            "Always set Port Status = ON after typing IP",
            "PDU Failed? Check port status FIRST!"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 52,
        "title": "Common Mistakes: Wrong Default Gateway",
        "script": (
            "The SECOND most common mistake: wrong default gateway on PCs. The default gateway "
            "must be the IP address of the ROUTER INTERFACE that is on the SAME network as the PC. "
            "Example: If your PC is 192.168.1.10 and your router port is 192.168.1.1, then "
            "the gateway is 192.168.1.1. NOT 192.168.2.1 or any other number! If you have two "
            "networks (like our Cebu university), each side has its OWN gateway. Library PCs use "
            "the router IP on the Library side. Science PCs use the router IP on the Science side. "
            "Wrong gateway = PC cannot reach other networks = Failed PDU!"
        ),
        "music": "Correction music - 'Right Track' style",
        "animation": "Gateway arrows pointing to correct router interface",
        "emotion": "Corrective - fixing the second biggest mistake",
        "bullets": [
            "#2 MISTAKE: Wrong Default Gateway!",
            "Gateway = Router IP on YOUR network",
            "PC: 192.168.1.10 > Gateway: 192.168.1.1",
            "NOT 192.168.2.1 or anything else!",
            "Multi-network: each side uses its own router IP",
            "Wrong gateway = cannot reach other networks",
            "Failed PDU across networks? Check gateway!"
        ],
        "bg": "FDEBD0"
    },

    {
        "num": 53,
        "title": "Common Mistakes: DHCP vs Static Confusion",
        "script": (
            "Mistake number 3: Leaving PCs on DHCP when there is no DHCP server! When you click "
            "Desktop > IP Configuration on a PC, the default is sometimes DHCP. DHCP means the "
            "PC waits for a server to give it an IP address automatically. But if you do not "
            "have a DHCP server in your network, the PC gets nothing! It will have IP 0.0.0.0 "
            "or a weird 169.254.x.x address. Always click Static and type the IP manually "
            "for our exercises. You can tell it is wrong if the IP field shows 0.0.0.0 or "
            "starts with 169.254."
        ),
        "music": "Alert music - 'Attention Please' style",
        "animation": "DHCP vs Static radio button highlight",
        "emotion": "Preventive - catch this before testing",
        "bullets": [
            "#3 MISTAKE: PC set to DHCP with no DHCP server",
            "DHCP = waiting for automatic IP (from server)",
            "No DHCP server = PC gets NO address!",
            "IP shows 0.0.0.0 or 169.254.x.x = wrong!",
            "ALWAYS select Static for our exercises",
            "Then manually type IP, Mask, Gateway",
            "Static = you control the address"
        ],
        "bg": "FFFFFF"
    },
    {
        "num": 54,
        "title": "Common Mistakes: Cable Types and Red Dots",
        "script": (
            "Mistake 4: Using the wrong cable type! Red dots on a cable mean the connection is "
            "NOT working. Straight-Through: PC to Switch, Switch to Router (DIFFERENT devices). "
            "Cross-Over: PC to PC, Switch to Switch, Router to Router (SIMILAR devices). "
            "If you see red dots, delete the cable (click Delete tool, then click the cable) "
            "and reconnect with the correct type. Or just use the Auto cable tool (lightning "
            "bolt with question mark) and it picks the right one. Green dots = good. Red = bad."
        ),
        "music": "Quick fix music - 'Easy Solution' style",
        "animation": "Red dots turning green after cable fix",
        "emotion": "Quick fix energy - easy to solve!",
        "bullets": [
            "#4 MISTAKE: Wrong cable type = red dots!",
            "Straight-Through: DIFFERENT devices",
            "Cross-Over: SIMILAR devices",
            "Red dots? Delete cable, use correct type",
            "Or use Auto cable tool (lightning+?)",
            "Green dots = physical connection works",
            "Red dots = fix the cable first!"
        ],
        "bg": "D6EAF8"
    },

    {
        "num": 55,
        "title": "Common Mistakes: Subnet Mask Mismatch",
        "script": (
            "Mistake 5: Subnet mask mismatch between devices! ALL devices on the same network "
            "must use the SAME subnet mask. If your router says 255.255.255.0 but a PC says "
            "255.255.0.0, they might not communicate properly. The subnet mask tells the device "
            "which part of the IP address is the network portion and which is the host portion. "
            "For our exercises, ALWAYS use 255.255.255.0 on every device in the same network. "
            "If even ONE device has a different mask, your PDU test might show Failed!"
        ),
        "music": "Consistency music - 'Match Up' style",
        "animation": "Subnet masks compared side by side",
        "emotion": "Emphasizing consistency - everything must match",
        "bullets": [
            "#5 MISTAKE: Different subnet masks on same network!",
            "ALL devices must use SAME mask",
            "Our standard: 255.255.255.0 everywhere",
            "Mask tells device: what is my network?",
            "Different masks = confusion = failed communication",
            "Even 1 device with wrong mask can cause failure",
            "Always double-check: 255.255.255.0"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 56,
        "title": "Quick Troubleshooting Checklist (Use Before Every PDU Test)",
        "script": (
            "Before you test with the Add Simple PDU mailbox icon, run through this checklist "
            "mentally. Check 1: Are all cables showing green dots? If not, fix cables first. "
            "Check 2: Is the router port status ON for every interface? Check 3: Are all PCs "
            "set to Static (not DHCP)? Check 4: Do all PCs have the correct IP, mask, and gateway? "
            "Check 5: Is the gateway the router IP on the SAME network as the PC? "
            "Check 6: Is the subnet mask the same on all devices (255.255.255.0)? "
            "If all 6 checks pass, your PDU test should show Successful!"
        ),
        "music": "Checklist music - 'Pre-Flight Check' style",
        "animation": "Six checkboxes being ticked off",
        "emotion": "Organized and thorough - be systematic",
        "bullets": [
            "CHECKLIST before PDU test:",
            "1. Green dots on all cables?",
            "2. Router port status ON?",
            "3. PCs set to Static?",
            "4. Correct IP + Mask + Gateway on all PCs?",
            "5. Gateway = Router IP on SAME network?",
            "6. Subnet mask matches everywhere?",
            "All 6 pass = Successful PDU test!"
        ],
        "bg": "FDEBD0"
    },

    # ===== WRAP-UP AND CLOSING (Slides 57-60) =====
    {
        "num": 57,
        "title": "Summary: IP Configuration Rules (No CLI Needed!)",
        "script": (
            "Let us summarize the IP configuration rules we learned. For PCs: Click PC > "
            "Desktop tab > IP Configuration > Static > type IP, Mask, Gateway. For Routers: "
            "Click Router > Config tab > click Interface > type IP, Mask > Port Status ON. "
            "For Servers and Printers: Click > Config tab > click Interface (FastEthernet0) > "
            "type IP, Mask, Gateway > Port Status ON. For Switches (management IP): Click > "
            "Config > VLAN1 > type IP, Mask. For testing: Use Add Simple PDU (envelope icon, "
            "right toolbar). Click source, click destination, check bottom panel. "
            "No terminal, no CLI, no commands needed!"
        ),
        "music": "Summary and recap music - 'Looking Back' style",
        "animation": "Summary cards for each device type",
        "emotion": "Confident recap - you know this now!",
        "bullets": [
            "PC: Desktop > IP Configuration > Static",
            "Router: Config > Interface > IP + Mask + ON",
            "Server/Printer: Config > Interface > IP + Mask + GW + ON",
            "Switch VLAN1: Config > VLAN1 > IP + Mask",
            "Test: Add Simple PDU (envelope icon, right toolbar)",
            "Click Source > Click Destination > Check result",
            "NO terminal, NO CLI, NO commands!"
        ],
        "bg": "FFFFFF"
    },
    {
        "num": 58,
        "title": "Summary: What We Built Today",
        "script": (
            "Look at everything we built today! Scenario 1: Manila school lab with 4 PCs, "
            "router, and switch. Scenario 2: Barangay health center with PCs, printer, and server. "
            "Scenario 3: Sari-sari store with wired and wireless devices. Scenario 4: University "
            "campus with TWO networks connected by one router. Scenario 5: Government office with "
            "public and private zones. Scenario 6: Pisonet with 6 customer PCs and a management "
            "server. You configured IPs on routers, PCs, servers, printers, and smartphones. "
            "You tested all of them with the Add Simple PDU mailbox icon. Amazing work!"
        ),
        "music": "Achievement celebration music - 'Victory Lap' style",
        "animation": "All 6 scenarios shown as mini diagrams",
        "emotion": "Proud and accomplished - look what you did!",
        "bullets": [
            "Scenario 1: Manila School (4 PCs + Router + Switch)",
            "Scenario 2: Barangay Health (PCs + Printer)",
            "Scenario 3: Sari-Sari Store (Wired + Wireless)",
            "Scenario 4: University (TWO networks + routing)",
            "Scenario 5: Government (Public + Private zones)",
            "Scenario 6: Pisonet (6 PCs + Server)",
            "All tested with Add Simple PDU!"
        ],
        "bg": "D6EAF8"
    },

    {
        "num": 59,
        "title": "Questions and Answers",
        "script": (
            "Now it is time for your questions! Do you have anything you want to ask? Maybe you "
            "are confused about where to click for IP Configuration. Maybe you want to know why "
            "the PDU test failed. Maybe you want to see me configure a device again. There are "
            "no silly questions! Every question helps you learn. Remember, the key steps are: "
            "place devices, cable them, set IPs using GUI tabs, and test with the mailbox icon. "
            "If you think of questions later, post in Cisco NetAcad forums. The community helps!"
        ),
        "music": "Open and welcoming music - 'Friendly Chat' style",
        "animation": "Question mark animations floating up",
        "emotion": "Open and welcoming - every question is valid",
        "bullets": [
            "Ask anything!",
            "Confused about IP Configuration steps?",
            "PDU test failed? Let us troubleshoot!",
            "Want to see a device configured again?",
            "No silly questions!",
            "Cisco NetAcad forums for later questions",
            "Community is always ready to help"
        ],
        "bg": "D5F5E3"
    },
    {
        "num": 60,
        "title": "Thank You and Happy Networking!",
        "script": (
            "Thank you so much for joining us today! You learned how to configure IP addresses "
            "on every type of device using ONLY the GUI - no terminal needed. You learned to test "
            "connectivity using the Add Simple PDU mailbox icon on the right toolbar. You built "
            "6 complete Philippine scenarios from scratch. Remember: practice is the best teacher. "
            "Open Packet Tracer, build networks, test with the mailbox icon, and have fun! "
            "Keep learning, keep practicing, and one day you will build real networks that help "
            "real people. Maraming salamat and happy networking, everyone!"
        ),
        "music": "Celebration and farewell - 'Happy Ending' by Bensound",
        "animation": "Confetti and thank you message with fade out",
        "emotion": "Grateful and celebratory - warm farewell",
        "bullets": [
            "Thank you for joining!",
            "GUI only - no terminal needed!",
            "Add Simple PDU = your testing tool!",
            "6 Philippine scenarios completed!",
            "Practice is the best teacher",
            "Open PT, build, test, have fun!",
            "Maraming salamat! Happy networking!"
        ],
        "bg": "FDEBD0"
    },
]



def generate_word_document():
    """Generate the Word document with the webinar script."""
    doc = Document()

    # Title
    title = doc.add_heading("Cisco Packet Tracer Webinar Script", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(
        "Complete Narration Script for 1-Hour+ Webinar\n"
        "Language: Simple English (explain like talking to a 5-year-old)\n"
        "Total Slides: 60 | Estimated Time: ~1.5 minutes per slide\n"
        "IP Configuration Method: GUI Only (Desktop/Config tabs)\n"
        "Testing Method: Add Simple PDU (Mailbox/Envelope icon, Right toolbar)\n"
        "NO terminal, NO CLI, NO Command Prompt used anywhere"
    )
    doc.add_paragraph("")  # spacer

    for slide in SLIDES:
        # Slide heading
        doc.add_heading(f"Slide {slide['num']}: {slide['title']}", level=1)

        # Script/Narration
        doc.add_heading("Script (What to Say):", level=2)
        p = doc.add_paragraph(slide["script"])
        p.style.font.size = Pt(12)

        # Music suggestion
        doc.add_heading("Background Music:", level=2)
        doc.add_paragraph(slide["music"])

        # Animation suggestion
        doc.add_heading("Animation/Transition:", level=2)
        doc.add_paragraph(slide["animation"])

        # Emotion/delivery
        doc.add_heading("Emotion/Delivery Style:", level=2)
        doc.add_paragraph(slide["emotion"])

        # Bullet points for reference
        doc.add_heading("Key Points on Slide:", level=2)
        for bullet in slide["bullets"]:
            doc.add_paragraph(bullet, style="List Bullet")

        # Separator
        doc.add_paragraph("_" * 60)
        doc.add_paragraph("")  # spacer

    output_path = "webinar/Cisco_Packet_Tracer_Webinar_Script.docx"
    doc.save(output_path)
    print(f"Word document saved: {output_path}")



def generate_powerpoint():
    """Generate the PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color map for backgrounds
    bg_colors = {
        "D6EAF8": PptRGB(0xD6, 0xEA, 0xF8),  # light blue
        "D5F5E3": PptRGB(0xD5, 0xF5, 0xE3),  # light green
        "FDEBD0": PptRGB(0xFD, 0xEB, 0xD0),  # light orange
        "FFFFFF": PptRGB(0xFF, 0xFF, 0xFF),  # white
    }

    for slide_data in SLIDES:
        # Use blank layout
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Set background color
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_colors.get(slide_data["bg"], PptRGB(0xFF, 0xFF, 0xFF))

        # Add title text box
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.3)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Slide {slide_data['num']}: {slide_data['title']}"
        p.font.size = PptPt(32)
        p.font.bold = True
        p.font.color.rgb = PptRGB(0x1A, 0x23, 0x7E)
        p.alignment = PP_ALIGN.LEFT

        # Add bullet points
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(11.5)
        height = Inches(5.2)
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, bullet in enumerate(slide_data["bullets"]):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = bullet
            p2.font.size = PptPt(20)
            p2.font.color.rgb = PptRGB(0x2C, 0x3E, 0x50)
            p2.space_after = PptPt(8)
            p2.level = 0

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = (
            f"SCRIPT: {slide_data['script']}\n\n"
            f"MUSIC: {slide_data['music']}\n"
            f"ANIMATION: {slide_data['animation']}\n"
            f"EMOTION: {slide_data['emotion']}"
        )

    output_path = "webinar/Cisco_Packet_Tracer_Presentation.pptx"
    prs.save(output_path)
    print(f"PowerPoint saved: {output_path}")



if __name__ == "__main__":
    print("Generating Cisco Packet Tracer Webinar materials...")
    print(f"Total slides: {len(SLIDES)}")
    generate_word_document()
    generate_powerpoint()
    print("Done! Both files are in the webinar/ folder.")
    print("KEY CHANGES in this version:")
    print("  - 60 slides (expanded from 40)")
    print("  - Detailed step-by-step IP placement for EVERY device")
    print("  - GUI ONLY (Desktop tab for PCs, Config tab for Routers/Servers/Printers)")
    print("  - NO terminal, NO CLI, NO Command Prompt anywhere")
    print("  - All testing uses Add Simple PDU (mailbox/envelope icon, right toolbar)")
    print("  - 6 Philippine scenarios with full IP plans and configurations")
